from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(0, 51, 102)
ACCENT_BLUE = RGBColor(0, 112, 192)
LIGHT_BLUE = RGBColor(68, 114, 196)
GREEN = RGBColor(112, 173, 71)
ORANGE = RGBColor(255, 124, 128)
GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(217, 217, 217)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add background rectangle
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Add accent bar
    accent = slide.shapes.add_shape(
        1,
        0, Inches(3), prs.slide_width, Inches(0.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_BLUE
    accent.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_items, is_two_column=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title bar
    title_bar = slide.shapes.add_shape(
        1,
        0, 0, prs.slide_width, Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Add content
    if is_two_column:
        col_width = Inches(5.5)
        left_col = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), col_width, Inches(5.5))
        right_col = slide.shapes.add_textbox(Inches(7), Inches(1.3), col_width, Inches(5.5))

        mid_point = len(content_items) // 2
        add_bullet_text(left_col.text_frame, content_items[:mid_point])
        add_bullet_text(right_col.text_frame, content_items[mid_point:])
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        add_bullet_text(content_box.text_frame, content_items)

def add_bullet_text(text_frame, items):
    text_frame.word_wrap = True
    for i, item in enumerate(items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        if isinstance(item, dict):
            p.text = item['text']
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 18))
            if item.get('color'):
                p.font.color.rgb = item['color']
        else:
            p.text = item
            p.font.size = Pt(18)

        p.font.color.rgb = p.font.color.rgb if hasattr(p.font.color, 'rgb') and p.font.color.rgb else GRAY

# Slide 1: Title
add_title_slide(
    "EDP Director Summary",
    "September 17-26, 2025 | PI 3.2 Planning & Infrastructure Migration"
)

# Slide 2: Executive Summary
add_content_slide(
    "Executive Summary",
    [
        "✓ Successfully completed PI 3.2 planning with full team alignment",
        "✓ Initiated critical infrastructure consolidation to NorthStar Snowflake account",
        "✓ Established business alignment strategy with Data Domain Councils",
        "✓ Secured Abacus partnership for CMS interoperability requirements",
        "⚠ Navigated repository security incident with minimal disruption",
        "",
        "Key Focus: Single Snowflake account, real-time data capability, business-driven roadmap"
    ]
)

# Slide 3: Critical Decisions - Infrastructure
add_content_slide(
    "Critical Decisions: Infrastructure Consolidation",
    [
        {"text": "Environment Migration Strategy", "size": 22, "color": ACCENT_BLUE},
        "• Migrate dev from Rising Sun to NorthStar immediately",
        "• Rename 'Franken' environment to 'Dev' (fast approach)",
        "• Rising Sun repurposed for POC/training with synthetic data",
        "",
        {"text": "Database Architecture", "size": 22, "color": ACCENT_BLUE},
        "• Separate raw databases: dev_raw_DB (ingestion) + dev_eng_raw_DB (engineering)",
        "• Zero-copy clone approach for fresh production-like data",
        "• Enables independent team workflows",
        "",
        {"text": "Timeline: End of Sprint 1 (mid-October)", "size": 20, "color": GREEN}
    ]
)

# Slide 4: Critical Decisions - Real-Time Data
add_content_slide(
    "Critical Decisions: Real-Time Data Architecture",
    [
        {"text": "Deferred Merge View Pattern (Selected)", "size": 22, "color": ACCENT_BLUE},
        "• Batch loading remains at 4-6 hour cycle",
        "• Real-time views provide seconds-level latency for critical tables",
        "• Existing working pattern repurposed from Kafka development",
        "• Target: <5 minute latency for customer service use cases",
        "",
        {"text": "Ingestion Simplification", "size": 22, "color": ACCENT_BLUE},
        "• Move from Kafka/MSK to S3 CSV file ingestion",
        "• Significant cost reduction with adequate performance",
        "• MSK retained only for true real-time requirements",
        "",
        {"text": "OneView app pipelines refactored for CSV-based streaming", "size": 18, "color": GRAY}
    ]
)

# Slide 5: Critical Decisions - Repository Strategy
add_content_slide(
    "Critical Decisions: Code Repository Split",
    [
        {"text": "Current State: Monolithic Repository", "size": 22, "color": ORANGE},
        "• Single repository causes release schedule conflicts",
        "• Cross-team dependencies slow delivery",
        "",
        {"text": "Target State: Three Independent Projects", "size": 22, "color": GREEN},
        "• EDP_Streaming: Real-time OneView pipelines",
        "• EDP_Data_Integrations: Extract projects",
        "• EDP_Data_Domains: Analytical platform transformations",
        "",
        {"text": "Benefits", "size": 22, "color": ACCENT_BLUE},
        "• Autonomous team delivery with independent release schedules",
        "• Reduced cross-team dependencies",
        "• OneView can iterate rapidly without impacting analytics"
    ]
)

# Slide 6: Key Accomplishments
add_content_slide(
    "Key Accomplishments",
    [
        {"text": "✓ PI 3.2 Planning Execution", "size": 20, "color": GREEN},
        "  All teams aligned on features, stories, dependencies for 10-week roadmap",
        "",
        {"text": "✓ Infrastructure Migration Progress", "size": 20, "color": GREEN},
        "  EDP_source_data migrated; EDP_Data_Domains 90% complete",
        "  Single Snowflake account reduces complexity and cost",
        "",
        {"text": "✓ Abacus Partnership Launch", "size": 20, "color": GREEN},
        "  CMS compliance (9115, 0057) - SOW signed, work begins October",
        "  Snowflake share for raw + silver foundational data marts",
        "",
        {"text": "✓ C4-DDD Architecture Framework", "size": 20, "color": GREEN},
        "  Executive diagrams (ECC), technical diagrams (ARB), detailed specs (Engineering)",
        "  Enables multi-level stakeholder communication"
    ]
)

# Slide 7: Escalation Items
add_content_slide(
    "Escalation Items & Resolution",
    [
        {"text": "⚠ Repository Security Incident - RESOLVED", "size": 22, "color": ORANGE},
        "• Repository flagged for PHI (PHI-like test data) and deleted without notice",
        "• Impact: Lost all branches except develop, all merge history",
        "• Resolution: Restored from backup within 24 hours, all environments operational",
        "• Note: Would have been production-down event if prod were live",
        "",
        {"text": "Recommendations", "size": 20, "color": ACCENT_BLUE},
        "  → Implement notification process before repository deletion",
        "  → Improve test data generation standards to avoid PHI-like patterns",
        "  → Establish clearer backup/restore procedures",
        "",
        {"text": "⚠ Snowflake AI Cortex Cost Spike", "size": 22, "color": ORANGE},
        "• Unexpected budget consumption detected",
        "• Action: Engage Snowflake AI expert for cost-effective POC guidance"
    ]
)

# Slide 8: Action Items - High Priority
add_content_slide(
    "Action Items: High Priority",
    [
        {"text": "Dan Brickey", "size": 20, "color": ACCENT_BLUE},
        "• Complete dbt testing in NorthStar Dev",
        "• Finalize Rising Sun decommissioning plan",
        "  Due: End of Sprint 1 (mid-October)",
        "",
        {"text": "Dan / Lindsay / Data Governance", "size": 20, "color": ACCENT_BLUE},
        "• Launch Membership and Product Data Domain Councils",
        "  Due: Early Q4 2025",
        "",
        {"text": "Data Engineering Teams", "size": 20, "color": ACCENT_BLUE},
        "• Split repositories: Streaming, Integrations, Domains",
        "• Enable independent release schedules",
        "  Due: Mid-PI 3.2"
    ]
)

# Slide 9: Action Items - Medium Priority
add_content_slide(
    "Action Items: Medium Priority",
    [
        {"text": "Provider 360 Modeling", "size": 20, "color": ACCENT_BLUE},
        "Owner: Dan / Nicole Bowen / Tina Day",
        "• Complete business modeling and metric naming alignment",
        "Due: October 2025",
        "",
        {"text": "OneView AI Proof of Concept", "size": 20, "color": ACCENT_BLUE},
        "Owner: Dan Brickey",
        "• Contract document harvesting design meetings",
        "• Target Q1 2026 POC, design complete by end Q4 2025",
        "",
        {"text": "Dev Environment Database Structure", "size": 20, "color": ACCENT_BLUE},
        "Owner: Engineering Admins",
        "• Create dual raw DB (ingestion + engineering with zero-copy clone)",
        "Due: Sprint 1"
    ]
)

# Slide 10: Strategic Progress Timeline
add_content_slide(
    "Strategic Progress: Environment Consolidation",
    [
        {"text": "Current State", "size": 22, "color": ORANGE},
        "• Development split across Rising Sun and NorthStar accounts",
        "• Franken environment in NorthStar used for testing",
        "",
        {"text": "In Progress", "size": 22, "color": LIGHT_BLUE},
        "• Migration to NorthStar (Franken→Dev) in final testing",
        "• Rising Sun decommissioning plan in development",
        "",
        {"text": "Target State (End of PI)", "size": 22, "color": GREEN},
        "• All dev, test, and production in single NorthStar account",
        "• Rising Sun repurposed for POC with synthetic data",
        "• Clean slate architecture without legacy technical debt"
    ]
)

# Slide 11: Strategic Progress - Real-Time Data
add_content_slide(
    "Strategic Progress: Real-Time Data Capability",
    [
        {"text": "Current State", "size": 22, "color": ORANGE},
        "• Batch loading twice daily (4-6 hour cycles)",
        "• MSK streaming partially implemented but expensive",
        "",
        {"text": "In Progress", "size": 22, "color": LIGHT_BLUE},
        "• Deferred merge view pattern implementation",
        "• CSV-based ingestion architecture development",
        "",
        {"text": "Target State", "size": 22, "color": GREEN},
        "• <5 minute latency for customer service critical tables",
        "• Cost-effective S3 CSV ingestion",
        "• MSK retired except for true real-time needs",
        "• Dynamic table layers for business rule application"
    ]
)

# Slide 12: Strategic Progress - Business Alignment
add_content_slide(
    "Strategic Progress: Business Alignment",
    [
        {"text": "Current State", "size": 22, "color": ORANGE},
        "• Provider Data Council established and active",
        "• Ad-hoc business engagement for other domains",
        "",
        {"text": "In Progress", "size": 22, "color": LIGHT_BLUE},
        "• Membership and Product Council formation",
        "• Provider 360 modeling alignment",
        "• C4-DDD framework development for stakeholder communication",
        "",
        {"text": "Target State", "size": 22, "color": GREEN},
        "• All major domains have active Data Councils",
        "• Business-driven roadmap prioritization",
        "• MDM solution proposal from Hakoda partner",
        "• Shared vocabulary through domain-driven design"
    ]
)

# Slide 13: Team Structure Evolution
add_content_slide(
    "Team Structure & Responsibilities",
    [
        {"text": "New Team Assignment: Business Alignment", "size": 22, "color": ACCENT_BLUE},
        "• Moved from individual contributor to solution architecture team",
        "• Focus on business domain alignment and data council leadership",
        "",
        {"text": "Key Responsibilities", "size": 20, "color": ACCENT_BLUE},
        "• Solution architecture oversight for environment and repository changes",
        "• Business domain modeling: Member, Product, Provider",
        "• Hakoda offshore team coordination and work assignment",
        "• Data Domain Council facilitation with business stakeholders",
        "",
        {"text": "Delegation Strategy", "size": 20, "color": GREEN},
        "• Distribute tangible technical work to team members",
        "• Spread guidance across multiple teams",
        "• Reserve capacity for business alignment activities"
    ]
)

# Slide 14: Risk & Dependency Management
add_content_slide(
    "Risk & Dependency Management",
    [
        {"text": "Resource Allocation Risks", "size": 20, "color": ORANGE},
        "• Architecture team capacity stretched across multiple initiatives",
        "• Balance needed between delegation and hands-on technical leadership",
        "",
        {"text": "Cross-Team Dependencies", "size": 20, "color": ORANGE},
        "• OneView real-time depends on CSV ingestion architecture completion",
        "• Repository splitting depends on successful NorthStar migration",
        "• Data Council success requires C4-DDD framework adoption",
        "",
        {"text": "Vendor & Partner Coordination", "size": 20, "color": ORANGE},
        "• Abacus timeline extends to summer 2026 (must align data model maturity)",
        "• Hakoda MDM proposal Q1 must align with Data Council outputs",
        "• Snowflake consultant availability critical for architecture validation"
    ]
)

# Slide 15: Next Steps & Milestones
add_content_slide(
    "Next Steps & Key Milestones",
    [
        {"text": "Sprint 1 Focus (Next 2 Weeks)", "size": 22, "color": ACCENT_BLUE},
        "✓ Complete NorthStar Dev environment testing",
        "✓ Finalize Test and Prod configuration migration",
        "✓ Begin repository split planning for OneView/Streaming",
        "✓ Schedule Membership/Product Data Council kickoffs",
        "",
        {"text": "Mid-PI Milestones", "size": 22, "color": ACCENT_BLUE},
        "✓ Rising Sun fully decommissioned",
        "✓ Repository split completed with independent release schedules",
        "✓ Real-time deferred merge view pattern implemented and tested",
        "✓ First Membership/Product Council meetings held",
        "",
        {"text": "Next Director Summary: October 7, 2025", "size": 20, "color": GREEN}
    ]
)

# Save presentation
output_path = r"C:\Users\danbr\github-danbrickey\edp-ai-expert-team\docs\presentations\EDP_Director_Summary_2025-09-17-26.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")